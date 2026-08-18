# -*- coding: utf-8 -*-
"""Áp dụng nâng cấp vào 8 deck: speaker notes + slide phân cách mục + slide số liệu."""
import os
from pptx import Presentation
from upgrade_decks import add_section_divider, add_stat_slide, move_slide, set_notes
from notes_data import NOTES

DECKS = {
    "c1": "CHUONG 1 - TONG QUAN GIAO TIEP TRONG KINH DOANH.pptx",
    "c2": "CHUONG 2 - KY NANG GIAO TIEP CHUYEN NGHIEP.pptx",
    "c3": "CHUONG 3 - GIAO TIEP TRONG TINH HUONG DAC THU.pptx",
    "c4": "CHUONG 4 - DAM PHAN TRONG KINH DOANH.pptx",
    "c5": "CHUONG 5 - SOAN THAO VA TRINH BAY VAN BAN.pptx",
    "th1": "THUC HANH BAI 1 - THE THUC VAN BAN.pptx",
    "th2": "THUC HANH BAI 2 - SOAN THAO VAN BAN HANH CHINH.pptx",
    "th3": "THUC HANH BAI 3 - SOAN THAO VAN BAN THUONG MAI.pptx",
}
LABEL = {"c1": "Chương 1", "c2": "Chương 2", "c3": "Chương 3", "c4": "Chương 4", "c5": "Chương 5",
         "th1": "Thực hành – Bài 1", "th2": "Thực hành – Bài 2", "th3": "Thực hành – Bài 3"}

# slide số liệu nổi bật: (deck, chèn sau slide chứa cụm, kicker, tiêu đề, [(số, nhãn, mô tả)], footer, ghi chú)
STATS = [
    ("c1", "Phương tiện giao tiếp: phi ngôn ngữ", "Con số phải nhớ", "Thông điệp cảm xúc được truyền đi thế nào?",
     [("7%", "TỪ NGỮ", "Nội dung lời nói — phần nhỏ nhất, dù ta thường đầu tư nhiều nhất vào đây"),
      ("38%", "GIỌNG NÓI", "Âm lượng, tốc độ, ngữ điệu, khoảng dừng"),
      ("55%", "NGÔN NGỮ CƠ THỂ", "Ánh mắt, nét mặt, cử chỉ, tư thế, khoảng cách")],
     "Nghiên cứu của Albert Mehrabian — chỉ áp dụng cho thông điệp mang tính CẢM XÚC, không phải mọi tình huống giao tiếp.",
     """CON SỐ MEHRABIAN (2 phút)
• Chiếu slide, để lớp nhìn 5 giây rồi mới nói.
• NÓI RÕ GIỚI HẠN: con số này chỉ đúng với thông điệp cảm xúc — đừng suy ra "nội dung chỉ chiếm 7%".
• Ý nghĩa thực hành: khi lời nói và cơ thể mâu thuẫn, người nghe TIN CƠ THỂ.
• Hỏi lớp: "Bạn nhận ra ai đó đang không vui dù họ nói 'em ổn' bằng cách nào?"''"""),

    ("c2", "Ấn tượng ban đầu", "Con số phải nhớ", "Quy tắc 4 × 20 — bốn cửa ải của ấn tượng ban đầu",
     [("20", "GIÂY", "Đối phương hình thành đánh giá tổng thể"),
      ("20", "BƯỚC CHÂN", "Dáng đi, tư thế được đọc từ xa"),
      ("20", "CENTIMET", "Ánh mắt và nụ cười trên gương mặt"),
      ("20", "TỪ", "Lời chào và giới thiệu đầu tiên")],
     "Ấn tượng ban đầu hình thành gần như tức thì và rất khó đảo ngược — nó phải được CHUẨN BỊ, không phó mặc cho may mắn.",
     """QUY TẮC 4×20 (2 phút)
• Slide này để sinh viên CHỤP LẠI.
• Đọc từng con số, mỗi con số cho một ví dụ 1 câu.
• Hỏi: "Trong 4 cửa ải này, bạn tự tin nhất và yếu nhất ở đâu?"''"""),

    ("c4", "vũ khí quan trọng nhất", "Con số phải nhớ", "Vì sao chuẩn bị lại quan trọng đến thế?",
     [("70%", "KẾT QUẢ", "được quyết định ngay từ giai đoạn chuẩn bị, trước khi hai bên ngồi vào bàn"),
      ("3", "MỨC MỤC TIÊU", "Lý tưởng – Kỳ vọng – Tối thiểu: phải viết ra giấy trước khi đàm phán"),
      ("1", "BATNA", "Phương án thay thế tốt nhất — nguồn sức mạnh thật sự trên bàn đàm phán")],
     "Câu để đời của chương: “Không chuẩn bị chính là chuẩn bị để nhượng bộ.”",
     """CON SỐ CHUẨN BỊ (2 phút)
• Nhấn 70% — đây là con số sinh viên hay bất ngờ.
• Liên hệ: bài role-play cuối buổi sẽ chấm phần chuẩn bị (mục tiêu 3 mức + BATNA viết ra giấy).
• Câu chốt đọc chậm: "Không chuẩn bị chính là chuẩn bị để nhượng bộ."''"""),

    ("c5", "Kỹ thuật trình bày", "Con số phải nhớ", "Bốn con số của một trang văn bản đúng chuẩn",
     [("A4", "KHỔ GIẤY", "210 × 297 mm — không dùng khổ Letter"),
      ("30–35", "MILIMET", "Lề trái, rộng nhất để đóng gáy lưu trữ"),
      ("13–14", "CỠ CHỮ", "Times New Roman, bộ mã Unicode, màu đen"),
      ("9", "THÀNH PHẦN", "Số thành phần thể thức bắt buộc theo NĐ 30/2020")],
     "Lề trên và dưới 20 – 25 mm • lề phải 15 – 20 mm • số trang đánh từ trang thứ hai, canh giữa theo lề trên.",
     """BỐN CON SỐ (2 phút)
• Đây là slide sinh viên phải THUỘC LÒNG — sẽ ra thi và dùng suốt 3 buổi thực hành.
• Đọc chậm từng con số, yêu cầu cả lớp ghi vào vở.
• Mẹo nhớ lề trái rộng nhất: để ĐÓNG GÁY.
• Kiểm tra ngay: gọi 2 sinh viên nhắc lại không nhìn slide."""),
]

# slide phân cách mục: (deck, chèn TRƯỚC slide chứa cụm, số hiệu, tiêu đề, phụ đề, ghi chú)
DIVIDERS = [
    ("c1", "Giao tiếp và giao tiếp trong kinh doanh", "Mục 1.1 – 1.2", "Bản chất và phương tiện giao tiếp",
     "Hiểu đúng bản chất trước khi luyện kỹ năng: giao tiếp là gì, diễn ra qua những khâu nào, bằng phương tiện gì.",
     "CHUYỂN MỤC (30 giây) — nói một câu dẫn rồi sang slide tiếp. Đây là mốc để lớp biết mình đang ở đâu trong bài."),
    ("c1", "Yếu tố ảnh hưởng", "Mục 1.4 – 1.5", "Yếu tố ảnh hưởng và nguyên tắc",
     "Điều gì làm hỏng một cuộc giao tiếp, và năm nguyên tắc giúp ta tránh được điều đó.",
     "CHUYỂN MỤC (30 giây) — nhắc: phần còn lại của chương là phần hay ra thi."),
    ("c2", "Ấn tượng ban đầu", "Mục 2.1", "Ấn tượng ban đầu và nghi thức xã giao",
     "Hai mươi giây đầu tiên quyết định phần lớn cách người khác nhìn nhận bạn.",
     "CHUYỂN MỤC (30 giây)."),
    ("c2", "5 bước chuẩn bị", "Mục 2.2", "Kỹ năng thuyết trình",
     "Từ chuẩn bị đến trình bày: làm sao để người nghe nhớ được điều bạn muốn nói.",
     "CHUYỂN MỤC (30 giây) — nhắc bài thuyết trình nhóm cuối kỳ chấm theo phần này."),
    ("c2", "Nghe khác lắng nghe", "Mục 2.3 – 2.4", "Lắng nghe, đặt câu hỏi và giao tiếp qua điện thoại",
     "Hai kỹ năng ít được dạy nhất nhưng quyết định nhất trong công việc hằng ngày.",
     "CHUYỂN MỤC (30 giây)."),
    ("c3", "Giao tiếp với cấp trên", "Mục 3.1", "Giao tiếp trong nội bộ tổ chức",
     "Với cấp trên, cấp dưới và đồng nghiệp — mỗi mối quan hệ một cách ứng xử.",
     "CHUYỂN MỤC (30 giây)."),
    ("c3", "Giao tiếp với khách hàng", "Mục 3.2", "Giao tiếp với bên ngoài tổ chức",
     "Khách hàng, đối tác, cơ quan nhà nước và truyền thông — bốn nhóm, bốn luật chơi.",
     "CHUYỂN MỤC (30 giây) — báo trước: phần quan trọng nhất chương nằm ở đây (quy trình LAST)."),
    ("c3", "Giao tiếp trên bàn tiệc", "Mục 3.3 – 3.4", "Bàn tiệc và môi trường đa văn hóa",
     "Nơi công việc vẫn tiếp diễn dù không ai nhắc đến công việc.",
     "CHUYỂN MỤC (30 giây) — đây là phần sinh viên hào hứng nhất, giữ nhịp vui nhưng vẫn chuẩn mực."),
    ("c4", "Đàm phán là gì", "Mục 4.1", "Khái niệm và các kiểu đàm phán",
     "Hiểu bản chất kép của đàm phán: vừa hợp tác vừa cạnh tranh.",
     "CHUYỂN MỤC (30 giây)."),
    ("c4", "Tiến trình đàm phán", "Mục 4.2", "Tiến trình đàm phán năm giai đoạn",
     "Bảy mươi phần trăm kết quả được quyết định trước khi hai bên ngồi vào bàn.",
     "CHUYỂN MỤC (30 giây) — nhấn con số 70% ngay tại slide này."),
    ("c4", "Kỹ năng nền tảng", "Mục 4.3", "Kỹ năng và chiêu trò trên bàn đàm phán",
     "Những gì cần rèn, và những gì cần nhận diện để không bị dẫn dắt.",
     "CHUYỂN MỤC (30 giây)."),
    ("c5", "Văn bản là gì", "Mục 5.1 – 5.2", "Khái niệm, phân loại và thể thức văn bản",
     "Nền tảng pháp lý và kỹ thuật: văn bản là gì và một trang văn bản đúng chuẩn trông thế nào.",
     "CHUYỂN MỤC (30 giây)."),
    ("c5", "Năm văn bản hành chính", "Mục 5.3", "Soạn thảo văn bản hành chính thông dụng",
     "Năm loại văn bản dùng hằng ngày trong mọi cơ quan, tổ chức.",
     "CHUYỂN MỤC (30 giây) — nhắc: bài thực hành số 2 làm đúng năm loại này."),
    ("c5", "Thư tín thương mại và báo giá", "Mục 5.4", "Soạn thảo văn bản thương mại",
     "Thư tín, báo giá và hợp đồng — bộ hồ sơ đưa một thương vụ đi đến đích.",
     "CHUYỂN MỤC (30 giây) — nhắc: bài thực hành số 3 làm đúng bộ hồ sơ này."),
]


def find_idx(prs, needle):
    for i, s in enumerate(prs.slides):
        t = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        if needle in t:
            return i
    return None


def apply_notes(prs, key):
    hit = 0
    table = NOTES.get(key, {})
    used = set()
    for i, s in enumerate(prs.slides):
        t = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        for needle, note in table.items():
            if needle in used:
                continue
            if needle in t:
                set_notes(s, note)
                used.add(needle)
                hit += 1
                break
    missing = [n for n in table if n not in used]
    return hit, missing


def main():
    for key, fn in DECKS.items():
        if not os.path.exists(fn):
            print("thiếu:", fn); continue
        prs = Presentation(fn)
        label = LABEL[key]

        # 1) chèn slide số liệu (chèn SAU slide đích)
        n_stat = 0
        for d, after, kicker, title, stats, footer, note in STATS:
            if d != key:
                continue
            idx = find_idx(prs, after)
            if idx is None:
                print(f"  ! {key}: không thấy slide '{after}' để chèn số liệu"); continue
            add_stat_slide(prs, kicker, title, stats, footer, label, note)
            move_slide(prs, len(prs.slides._sldIdLst) - 1, idx + 1)
            n_stat += 1

        # 2) chèn slide phân cách (chèn TRƯỚC slide đích)
        n_div = 0
        for d, before, num, title, sub, note in DIVIDERS:
            if d != key:
                continue
            idx = find_idx(prs, before)
            if idx is None:
                print(f"  ! {key}: không thấy slide '{before}' để chèn phân cách"); continue
            add_section_divider(prs, num, title, sub, label, note)
            move_slide(prs, len(prs.slides._sldIdLst) - 1, idx)
            n_div += 1

        # 3) speaker notes
        n_note, missing = apply_notes(prs, key)
        prs.save(fn)
        msg = f"{key}: +{n_div} slide phân cách, +{n_stat} slide số liệu, {n_note} slide có ghi chú giảng bài"
        if missing:
            msg += f"  (chưa khớp: {len(missing)})"
        print(msg)


if __name__ == "__main__":
    main()
