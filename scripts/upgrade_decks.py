# -*- coding: utf-8 -*-
"""Nâng cấp 8 bộ bài giảng EC1103 theo nguyên tắc từ các skill PPT:
1) Speaker notes dạng "tín hiệu nhắc" cho từng slide  -> dùng Presenter View khi giảng
2) Slide phân cách mục (section divider) trước mỗi mục lớn
3) Slide số liệu nổi bật (stat highlight) cho các con số phải thuộc
"""
import copy, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CORAL = RGBColor(0xDC, 0x75, 0x6A)
RUST = RGBColor(0xAC, 0x4D, 0x33)
BLUSH = RGBColor(0xFB, 0xCE, 0xC9)
BLUSH_SOFT = RGBColor(0xFD, 0xF1, 0xEF)
CREAM = RGBColor(0xFD, 0xFB, 0xF8)
INK = RGBColor(0x3A, 0x2B, 0x28)
GRAY = RGBColor(0x8A, 0x7A, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD = "Cambria"
BODY = "Calibri"
SW, SH = 13.333, 7.5


def slide_text(s):
    return " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ---------- dựng slide mới ----------
def blank_layout(prs):
    for lay in prs.slide_layouts:
        if not lay.placeholders._element.findall(
                '{http://schemas.openxmlformats.org/presentationml/2006/main}sp'):
            return lay
    return prs.slide_layouts[-1]


def move_slide(prs, from_idx, to_idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[from_idx])
    xml_slides.insert(to_idx, slides[from_idx])


def chrome(slide, page_label, chapter_label):
    slide.shapes.add_picture("figs/logo-mark.png", Inches(0.55), Inches(7.0), Inches(0.34), Inches(0.34)) \
        if os.path.exists("figs/logo-mark.png") else None
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.02), Inches(8.5), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Je m'appelle Huong  •  GV. Đỗ Thùy Hương  •  EC1103"
    r.font.size = Pt(9); r.font.name = BODY; r.font.color.rgb = GRAY
    tb2 = slide.shapes.add_textbox(Inches(10.2), Inches(7.02), Inches(2.1), Inches(0.32))
    p2 = tb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = chapter_label
    r2.font.size = Pt(9); r2.font.name = BODY; r2.font.color.rgb = GRAY


def add_section_divider(prs, number, title, subtitle, chapter_label, notes=""):
    s = prs.slides.add_slide(blank_layout(prs))
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background(); bg.shadow.inherit = False
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.3), Inches(SW), Inches(3.0))
    band.fill.solid(); band.fill.fore_color.rgb = BLUSH_SOFT
    band.line.fill.background(); band.shadow.inherit = False
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), Inches(SH))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background(); bar.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(1.4), Inches(2.6), Inches(10.6), Inches(0.45))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = number.upper()
    r.font.size = Pt(15); r.font.bold = True; r.font.name = BODY
    r.font.color.rgb = CORAL

    tb = s.shapes.add_textbox(Inches(1.4), Inches(3.05), Inches(10.5), Inches(1.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    size = 38 if len(title) <= 34 else (33 if len(title) <= 46 else 29)
    r.font.size = Pt(size); r.font.bold = True; r.font.name = HEAD; r.font.color.rgb = RUST

    tb = s.shapes.add_textbox(Inches(1.4), Inches(4.62), Inches(10.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = subtitle
    r.font.size = Pt(14.5); r.font.italic = True; r.font.name = BODY; r.font.color.rgb = GRAY

    chrome(s, "", chapter_label)
    if notes:
        set_notes(s, notes)
    return s


def add_stat_slide(prs, kicker, title, stats, footer, chapter_label, notes=""):
    """stats: list of (số lớn, đơn vị/nhãn, mô tả ngắn)"""
    s = prs.slides.add_slide(blank_layout(prs))
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    bg.fill.solid(); bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background(); bg.shadow.inherit = False

    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = kicker.upper()
    r.font.size = Pt(11); r.font.bold = True; r.font.name = BODY; r.font.color.rgb = CORAL

    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.66), Inches(12.2), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(27); r.font.bold = True; r.font.name = HEAD; r.font.color.rgb = INK

    n = len(stats)
    gap = 0.26
    w = (12.25 - gap * (n - 1)) / n
    for i, (big, unit, desc) in enumerate(stats):
        x = 0.55 + i * (w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(w), Inches(3.75))
        card.adjustments[0] = 0.07
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLUSH; card.line.width = Pt(1.25)
        card.shadow.inherit = False
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.3); tf.margin_bottom = Inches(0.2)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = big
        r.font.size = Pt(58); r.font.bold = True; r.font.name = HEAD; r.font.color.rgb = CORAL
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = unit
        r2.font.size = Pt(14); r2.font.bold = True; r2.font.name = BODY; r2.font.color.rgb = RUST
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(12)
        r3 = p3.add_run(); r3.text = desc
        r3.font.size = Pt(12.5); r3.font.name = BODY; r3.font.color.rgb = INK

    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.95), Inches(12.25), Inches(0.75))
    box.adjustments[0] = 0.12
    box.fill.solid(); box.fill.fore_color.rgb = BLUSH_SOFT
    box.line.color.rgb = BLUSH; box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = footer
    r.font.size = Pt(13); r.font.name = BODY; r.font.color.rgb = INK

    chrome(s, "", chapter_label)
    if notes:
        set_notes(s, notes)
    return s


def renumber(prs):
    """Đánh lại số trang ở góc phải dưới (bỏ qua slide bìa)."""
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE.OVAL or (sh.has_text_frame and sh.text_frame.text.strip().isdigit()
                                                   and sh.top and sh.top > Inches(6.9)):
                if sh.has_text_frame and sh.text_frame.text.strip().isdigit():
                    for para in sh.text_frame.paragraphs:
                        for r in para.runs:
                            r.text = str(i)
