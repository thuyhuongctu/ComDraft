# -*- coding: utf-8 -*-
"""Dựng lại toàn bộ tám bộ slide từ đầu, một lệnh.

    python3 scripts/dung_slide.py           # dựng ra thư mục tạm rồi đối chiếu
    python3 scripts/dung_slide.py --ghi      # dựng xong thì ghi đè slides/ và practice/

CLAUDE.md dặn "sửa ở trình sinh rồi chạy lại, đừng sửa tay tệp kết quả". Nhưng
trước đây không ai chạy lại được: build_decks.js chết ngay slide đầu vì thiếu
gói pptxgenjs và vì design.js trỏ vào hai tệp logo không tồn tại, add_images.py
đọc thư mục "figs" trong khi repo tên là "figures", và không có bước nào đưa kết
quả về đúng tên trong slides/ với practice/. Tám bộ slide vì thế thành tệp mồ
côi — sửa được nhưng không dựng lại được.

Ba chặng phải chạy đúng thứ tự, mỗi chặng ghi đè tệp của chặng trước:

    build_decks.js + build_ch5_practice.js   dựng nội dung gốc
    apply_upgrade.py                          thêm ghi chú giảng bài, slide phân cách, slide số liệu
    add_images.py                             thay nội dung vài slide bằng hình minh họa

Mặc định KHÔNG ghi đè: dựng ra thư mục tạm rồi so số slide, số chữ và số ghi chú
với bộ đang dùng. Muốn thay thật thì thêm --ghi. Sai một chặng là mất cả bộ, mà
bộ đang dùng là bản cô Hương đã duyệt.
"""
import os
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation

GOC = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SCRIPTS = os.path.join(GOC, "scripts")

# Trình sinh đặt tên theo kiểu tài liệu phát cho lớp; repo lưu theo kiểu đường
# dẫn web. Thiếu đúng bảng ánh xạ này là chỗ dây chuyền đứt trước đây.
DAT_TEN = [
    ("CHUONG 1 - TONG QUAN GIAO TIEP TRONG KINH DOANH.pptx", "slides/01-tong-quan-giao-tiep.pptx"),
    ("CHUONG 2 - KY NANG GIAO TIEP CHUYEN NGHIEP.pptx", "slides/02-ky-nang-chuyen-nghiep.pptx"),
    ("CHUONG 3 - GIAO TIEP TRONG TINH HUONG DAC THU.pptx", "slides/03-tinh-huong-dac-thu.pptx"),
    ("CHUONG 4 - DAM PHAN TRONG KINH DOANH.pptx", "slides/04-dam-phan.pptx"),
    ("CHUONG 5 - SOAN THAO VA TRINH BAY VAN BAN.pptx", "slides/05-soan-thao-van-ban.pptx"),
    ("THUC HANH BAI 1 - THE THUC VAN BAN.pptx", "practice/bai-1-the-thuc.pptx"),
    ("THUC HANH BAI 2 - SOAN THAO VAN BAN HANH CHINH.pptx", "practice/bai-2-hanh-chinh.pptx"),
    ("THUC HANH BAI 3 - SOAN THAO VAN BAN THUONG MAI.pptx", "practice/bai-3-thuong-mai.pptx"),
]


def dac_diem(duong):
    """Ba con số đủ để biết hai bộ slide có cùng nội dung hay không."""
    p = Presentation(duong)
    chu = sum(len(sh.text_frame.text.split())
              for s in p.slides for sh in s.shapes if sh.has_text_frame)
    ghi = sum(1 for s in p.slides
              if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())
    return len(p.slides), chu, ghi


def chay(lenh, tai):
    moi = dict(os.environ)
    # apply_upgrade.py nạp upgrade_decks và notes_data nằm cạnh nó, nhưng lại
    # phải chạy ở thư mục chứa deck. Nên đưa scripts/ vào đường tìm mô-đun.
    moi["PYTHONPATH"] = SCRIPTS + os.pathsep + moi.get("PYTHONPATH", "")
    kq = subprocess.run(lenh, cwd=tai, env=moi, capture_output=True, text=True)
    if kq.returncode:
        sys.stderr.write(kq.stdout + kq.stderr)
        raise SystemExit("hỏng ở: " + " ".join(lenh))
    return kq.stdout


def main():
    ghi_de = "--ghi" in sys.argv
    tai = tempfile.mkdtemp(prefix="comdraft-slide-")
    try:
        for ten in ("build_decks.js", "build_ch5_practice.js"):
            chay(["node", os.path.join(SCRIPTS, ten)], tai)
        chay(["python3", os.path.join(SCRIPTS, "apply_upgrade.py")], tai)
        chay(["python3", os.path.join(SCRIPTS, "add_images.py")], tai)

        lech = 0
        print("%-32s %-20s %-20s" % ("", "DỰNG LẠI", "ĐANG DÙNG"))
        for nguon, dich in DAT_TEN:
            moi = os.path.join(tai, nguon)
            cu = os.path.join(GOC, dich)
            a = dac_diem(moi)
            b = dac_diem(cu) if os.path.exists(cu) else None
            hop = a == b
            lech += not hop
            print("%-32s %2d sl %4d chữ %2d ghi | %s  %s" % (
                os.path.basename(dich), a[0], a[1], a[2],
                ("%2d sl %4d chữ %2d ghi" % b) if b else "     (chưa có)      ",
                "khớp" if hop else "LỆCH"))
            if ghi_de:
                shutil.copyfile(moi, cu)

        if ghi_de:
            print("\nĐã ghi đè %d tệp." % len(DAT_TEN))
        elif lech:
            print("\n%d bộ lệch so với bản đang dùng. Xem kỹ rồi mới chạy --ghi." % lech)
        else:
            print("\nDựng lại khớp hoàn toàn bản đang dùng. Không cần ghi đè.")
        return 1 if (lech and not ghi_de) else 0
    finally:
        shutil.rmtree(tai, ignore_errors=True)


if __name__ == "__main__":
    raise SystemExit(main())
