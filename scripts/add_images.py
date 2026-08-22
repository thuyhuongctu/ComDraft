# -*- coding: utf-8 -*-
"""Chèn ảnh minh họa và ảnh nhân vật vào các bộ bài giảng EC1103 (hệ trực tiếp)."""
import glob, os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Trước đây là "figs" — một thư mục không có trong repo, nên mọi lần chèn hình
# đều lặng lẽ hỏng. Thư mục thật tên là "figures"; neo tuyệt đối theo vị trí
# tệp này để chạy từ thư mục dựng tạm nào cũng đúng.
F = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "figures")
CORAL = RGBColor(0xDC, 0x75, 0x6A)
RUST = RGBColor(0xAC, 0x4D, 0x33)
BLUSH = RGBColor(0xFB, 0xCE, 0xC9)
GRAY = RGBColor(0x8A, 0x7A, 0x76)
BODY = "Calibri"
SW, SH = 13.333, 7.5


def txt(slide, s):
    return " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def find_slide(prs, needle):
    for i, s in enumerate(prs.slides):
        if needle in txt(prs, s) if False else needle in " ".join(
                sh.text_frame.text for sh in s.shapes if sh.has_text_frame):
            return i, s
    return None, None


def clear_body(slide, keep_top=1.5):
    """Xóa các shape nội dung (giữ tiêu đề, footer, logo, số trang)."""
    for sh in list(slide.shapes):
        top = Emu(sh.top).inches
        if keep_top <= top < 6.85:
            sh._element.getparent().remove(sh._element)


def add_figure(slide, img, y=1.58, max_h=5.15, max_w=12.3, caption=None):
    from PIL import Image
    iw, ih = Image.open(img).size
    ar = iw / ih
    h = max_h
    w = h * ar
    if w > max_w:
        w = max_w
        h = w / ar
    x = (SW - w) / 2
    if caption:
        h -= 0.32
        w = h * ar
        if w > max_w:
            w = max_w; h = w / ar
        x = (SW - w) / 2
    # Hình bẹt hơn khung slide thì co theo bề rộng, cao ra ít hơn chỗ trống —
    # dán sát mép trên sẽ để hụt một mảng trắng ở đáy. Đặt vào giữa vùng thân
    # slide cho cân, phần trống chia đều trên dưới.
    vung = max_h + (0.32 if caption else 0)
    y = y + max(0.0, (vung - (h + (0.32 if caption else 0))) / 2)
    slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(w), Inches(h))
    if caption:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(y + h + 0.08), Inches(12.15), Inches(0.34))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.size = Pt(11.5); r.font.name = BODY; r.font.italic = True; r.font.color.rgb = GRAY
    return w, h


def add_persona_cover(slide, img=f"{F}/co-huong-doc.png", x=9.45, y=1.15, h=5.6):
    """Ảnh nhân vật ở mép phải slide bìa."""
    from PIL import Image
    iw, ih = Image.open(img).size
    w = h * iw / ih
    slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(w), Inches(h))


def add_persona_round(slide, x, y, d=1.5, img=f"{F}/co-huong-tron.png"):
    slide.shapes.add_picture(img, Inches(x), Inches(y), Inches(d), Inches(d))


def note_box(slide, text, y=6.15, h=0.62, color=RUST):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(12.25), Inches(h))
    sh.adjustments[0] = 0.12
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xFD, 0xF1, 0xEF)
    sh.line.color.rgb = BLUSH; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(12.5); r.font.name = BODY; r.font.color.rgb = color


# ---- bản đồ: file deck -> danh sách (từ khóa tìm slide, hình, chú thích) ----
PLAN = {
    "CHUONG 1 - TONG QUAN GIAO TIEP TRONG KINH DOANH.pptx": [
        ("Mô hình quá trình giao tiếp", f"{F}/c1-mo-hinh-giao-tiep-nt.png",
         "Giao tiếp là quá trình hai chiều — muốn sửa một cuộc giao tiếp thất bại, hãy dò lại từng khâu."),
        ("Các hình thức giao tiếp", f"{F}/c1-hinh-thuc-giao-tiep-nt.png",
         "Bốn cặp này không loại trừ nhau — một cuộc giao tiếp nằm đâu đó trên cả bốn trục cùng lúc."),
        ("Yếu tố ảnh hưởng", f"{F}/c1-yeu-to-anh-huong-nt.png",
         "Hỏng ở yếu tố nào thì sửa đúng yếu tố ấy, đừng đổ hết cho “nói chưa khéo”."),
    ],
    "CHUONG 2 - KY NANG GIAO TIEP CHUYEN NGHIEP.pptx": [
        ("Ấn tượng ban đầu", f"{F}/c2-quy-tac-4x20-nt.png", None),
        ("Nghe khác lắng nghe", f"{F}/c2-5-muc-lang-nghe-nt.png", None),
    ],
    "CHUONG 3 - GIAO TIEP TRONG TINH HUONG DAC THU.pptx": [
        ("quy trình LAST", f"{F}/c3-quy-trinh-last-nt.png", None),
    ],
    "CHUONG 4 - DAM PHAN TRONG KINH DOANH.pptx": [
        ("Tiến trình đàm phán", f"{F}/c4-tien-trinh-dam-phan-nt.png", None),
        ("vũ khí quan trọng nhất", f"{F}/c4-batna-zopa-nt.png",
         "ZOPA hẹp hay rộng phụ thuộc vào giới hạn thật của hai bên — chuẩn bị kỹ để biết mình đang ở đâu trên trục này."),
    ],
    "CHUONG 5 - SOAN THAO VA TRINH BAY VAN BAN.pptx": [
        ("Chín thành phần thể thức", f"{F}/c5-the-thuc-a4-nt.png", None),
        ("Năm văn bản hành chính thông dụng", f"{F}/c5-chuoi-van-ban-nt.png",
         "Chuỗi văn bản của một thương vụ — sinh viên sẽ soạn lại đúng chuỗi này trong phần thực hành."),
    ],
    "THUC HANH BAI 1 - THE THUC VAN BAN.pptx": [
        ("Trình bày các thành phần ở đầu văn bản", f"{F}/c5-the-thuc-a4-nt.png", None),
    ],
    "THUC HANH BAI 2 - SOAN THAO VAN BAN HANH CHINH.pptx": [
        ("Soạn thảo Quyết định", f"{F}/c5-chuoi-van-ban-nt.png",
         "Bài nộp số 2 nằm ở các bước 1 – 3 của chuỗi; bước 4 – 7 là nội dung Bài thực hành số 3."),
    ],
    "THUC HANH BAI 3 - SOAN THAO VAN BAN THUONG MAI.pptx": [
        ("Hợp đồng — các điều khoản phải có", f"{F}/c5-chuoi-van-ban-nt.png", None),
    ],
}


# Slide mục lục và slide mục tiêu liệt kê tên các mục của chương, nên chúng
# chứa gần như mọi từ khóa dùng để tìm slide đích. Nếu dò theo toàn bộ chữ trên
# slide thì hình minh họa sẽ rơi vào đúng hai slide này và xóa mất nội dung —
# lỗi đã xảy ra với mục lục Chương 4. Vì vậy chỉ dò trong vùng tiêu đề, và bỏ
# qua hẳn hai slide khung này.
SLIDE_KHUNG = ("Chúng ta sẽ đi qua", "sinh viên có thể")


def tieu_de_slide(s):
    """Chữ nằm ở vùng tiêu đề (mép trên dưới 1,6 inch)."""
    return " ".join(
        sh.text_frame.text for sh in s.shapes
        if sh.has_text_frame and sh.top is not None and sh.top < Inches(1.6)
    )


def la_slide_khung(s):
    t = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
    return any(k in t for k in SLIDE_KHUNG)


def process(fn, items):
    prs = Presentation(fn)
    n_fig = 0
    # 1) ảnh nhân vật trên slide bìa
    cover = prs.slides[0]
    add_persona_cover(cover)
    # 2) ảnh chân dung nhỏ ở slide "Hoạt động nhóm" / "Bài tập"
    for s in prs.slides:
        t = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        if "NHIỆM VỤ CỦA NHÓM" in t:
            add_persona_round(s, x=11.35, y=0.42, d=1.05)
            break
    # 3) chèn hình minh họa: thay nội dung slide đích bằng hình
    for needle, img, cap in items:
        idx, s = None, None
        for i, sl in enumerate(prs.slides):
            if la_slide_khung(sl):
                continue
            if needle in tieu_de_slide(sl):
                idx, s = i, sl
                break
        if s is None:
            print(f"  ! không tìm thấy slide có tiêu đề chứa: {needle}")
            continue
        clear_body(s)
        cap_y = 1.62
        add_figure(s, img, y=cap_y, caption=cap)
        n_fig += 1
    prs.save(fn)
    print(f"{fn}: {n_fig} hình minh họa + ảnh nhân vật")


if __name__ == "__main__":
    for fn, items in PLAN.items():
        if os.path.exists(fn):
            process(fn, items)
        else:
            print("thiếu file:", fn)
